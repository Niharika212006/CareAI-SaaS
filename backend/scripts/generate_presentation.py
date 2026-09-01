"""CareAI Final Year Project Presentation Generator.

Generates a professional, academically structured 22-slide PowerPoint (.pptx)
presentation for the CareAI Healthcare SaaS platform.
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# -----------------------------------------------------------------------------
# Color Palette & Typography
# -----------------------------------------------------------------------------
NAVY_900 = RGBColor(15, 23, 42)      # #0F172A - Main titles, dark headings
SLATE_800 = RGBColor(30, 41, 59)     # #1E293B - Card titles
SLATE_600 = RGBColor(71, 85, 105)    # #475569 - Body text
SLATE_500 = RGBColor(100, 116, 139)  # #64748B - Muted subtitles
SLATE_300 = RGBColor(203, 213, 225)  # #CBD5E1 - Borders
SLATE_100 = RGBColor(241, 245, 249)  # #F1F5F9 - Card fills
BG_LIGHT = RGBColor(248, 250, 252)   # #F8FAFC - Slide background
WHITE = RGBColor(255, 255, 255)      # #FFFFFF - Pure white

PRIMARY_BLUE = RGBColor(2, 132, 199) # #0284C7 - Primary accent (Cyan/Blue)
PRIMARY_DARK = RGBColor(3, 105, 161) # #0369A1
ACCENT_INDIGO = RGBColor(79, 70, 229)# #4F46E5 - AI & Tech accent
SUCCESS_GREEN = RGBColor(16, 185, 129)# #10B981 - Success & Verification
SUCCESS_BG = RGBColor(236, 253, 245) # #ECFDF5
WARNING_AMBER = RGBColor(217, 119, 6)# #D97706 - Warnings / Limits
WARNING_BG = RGBColor(254, 243, 199) # #FEF3C7
PURPLE_600 = RGBColor(147, 51, 234)  # #9333EA - Specialized roles
PURPLE_BG = RGBColor(250, 245, 255)  # #FAF5FF

FONT_HEADING = "Segoe UI"
FONT_BODY = "Segoe UI"

def create_presentation(output_filepath):
    prs = Presentation()
    # 16:9 Widescreen (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide_header(slide, title_text, category_text, slide_num):
        """Adds a standardized top banner and title block to content slides."""
        # Top category badge
        badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.35))
        tf_b = badge_box.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        p_b = tf_b.paragraphs[0]
        p_b.text = f"CAREAI PLATFORM  |  {category_text.upper()}"
        p_b.font.name = FONT_HEADING
        p_b.font.size = Pt(9.5)
        p_b.font.bold = True
        p_b.font.color.rgb = PRIMARY_BLUE

        # Slide Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(10), Inches(0.6))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = FONT_HEADING
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_900

        # Slide Number (Top Right)
        num_box = slide.shapes.add_textbox(Inches(11.0), Inches(0.45), Inches(1.5), Inches(0.4))
        tf_n = num_box.text_frame
        tf_n.margin_left = tf_n.margin_top = tf_n.margin_right = tf_n.margin_bottom = 0
        p_n = tf_n.paragraphs[0]
        p_n.alignment = PP_ALIGN.RIGHT
        p_n.text = f"{slide_num:02d} / 22"
        p_n.font.name = FONT_HEADING
        p_n.font.size = Pt(11)
        p_n.font.bold = True
        p_n.font.color.rgb = SLATE_500

        # Top separator line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = SLATE_300
        line.line.color.rgb = SLATE_300

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.3))
        tf_f = footer_box.text_frame
        tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = 0
        p_f = tf_f.paragraphs[0]
        p_f.text = "CareAI: AI-Powered Healthcare Management SaaS Platform  •  Final Year Project Evaluation"
        p_f.font.name = FONT_BODY
        p_f.font.size = Pt(8.5)
        p_f.font.color.rgb = SLATE_500

    def add_card(slide, left, top, width, height, bg_color=WHITE, border_color=SLATE_300):
        """Creates a rectangular card shape with solid fill and outline."""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    
    # Background
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY_900
    bg.line.fill.background()

    # Brand Badge
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(3.4), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PRIMARY_DARK
    badge.line.fill.background()
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "FINAL YEAR PROJECT PRESENTATION"
    p_b.font.name = FONT_HEADING
    p_b.font.size = Pt(10)
    p_b.font.bold = True
    p_b.font.color.rgb = WHITE
    p_b.alignment = PP_ALIGN.CENTER

    # Main Project Title
    t_box = s1.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(11), Inches(2.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "CareAI"
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Healthcare Management SaaS Platform"
    p2.font.name = FONT_HEADING
    p2.font.size = Pt(22)
    p2.font.color.rgb = PRIMARY_BLUE
    p2.space_before = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = "A Multi-Role, End-to-End Clinical Lifecycle Architecture with Role-Based Access Control & Gemini AI Clinical Safety Guardrails"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(13)
    p3.font.color.rgb = SLATE_300
    p3.space_before = Pt(12)

    # Feature Highlights Badges
    highlights = [
        ("5 Core User Roles", "Patient, Doctor, Admin, Lab Technician, Pharmacy Staff"),
        ("165/165 Automated Tests", "Comprehensive Pytest Backend Coverage with 100% Pass Rate"),
        ("Dual Cloud Deployment", "Vercel (React Frontend) + Render (FastAPI + Managed PostgreSQL)"),
        ("Google Gemini AI", "Safety Triage, Interaction Analysis & Graceful Fallback (HTTP 503)")
    ]
    for i, (hl_title, hl_sub) in enumerate(highlights):
        c_x = Inches(1.2 + (i % 2) * 5.6)
        c_y = Inches(4.5 + (i // 2) * 1.1)
        c_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_x, c_y, Inches(5.3), Inches(0.95))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        c_box.line.color.rgb = RGBColor(51, 65, 85)
        tf_c = c_box.text_frame
        tf_c.word_wrap = True
        p_ct = tf_c.paragraphs[0]
        p_ct.text = f"✔ {hl_title}"
        p_ct.font.name = FONT_HEADING
        p_ct.font.size = Pt(12)
        p_ct.font.bold = True
        p_ct.font.color.rgb = SUCCESS_GREEN
        p_cs = tf_c.add_paragraph()
        p_cs.text = hl_sub
        p_cs.font.name = FONT_BODY
        p_cs.font.size = Pt(10)
        p_cs.font.color.rgb = SLATE_300

    # =========================================================================
    # SLIDE 2: Project Overview
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_header(s2, "Executive Project Overview", "System Introduction", 2)

    overview_cards = [
        ("Centralized Clinical Hub", "Unifies isolated hospital departments into a single reactive platform, eliminating paper trails, physical handoffs, and disjointed records across care units.", PRIMARY_BLUE),
        ("5-Role Healthcare Ecosystem", "Provides tailor-made dashboards, permission guards, and dedicated operational tools for Patients, Doctors, Administrators, Lab Technicians, and Pharmacy Staff.", ACCENT_INDIGO),
        ("End-to-End Care Lifecycle", "Manages the entire patient journey: appointment booking, doctor consultations, digital prescriptions, lab sample chain-of-custody, and pharmacy dispensation.", SUCCESS_GREEN),
        ("Responsible AI Assistance", "Integrates Google Gemini 1.5 Flash for drug-drug interaction screening, medical document analysis, and health explanations with deterministic safety filters.", PURPLE_600)
    ]

    for i, (title, desc, color) in enumerate(overview_cards):
        x = Inches(0.8 + (i % 2) * 6.0)
        y = Inches(1.7 + (i // 2) * 2.5)
        add_card(s2, x, y, Inches(5.7), Inches(2.3), WHITE, SLATE_300)
        
        strip = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.2), Inches(5.3), Inches(0.45))
        strip.fill.solid()
        strip.fill.fore_color.rgb = SLATE_100
        strip.line.fill.background()
        tf_st = strip.text_frame
        p_st = tf_st.paragraphs[0]
        p_st.text = f"❖  {title}"
        p_st.font.name = FONT_HEADING
        p_st.font.size = Pt(13)
        p_st.font.bold = True
        p_st.font.color.rgb = color

        desc_box = s2.shapes.add_textbox(x + Inches(0.2), y + Inches(0.75), Inches(5.3), Inches(1.4))
        tf_d = desc_box.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = SLATE_600

    # =========================================================================
    # SLIDE 3: Problem Statement
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_header(s3, "Problem Statement & Clinical Motivation", "Clinical Challenges", 3)

    problems = [
        ("Fragmented Medical Workflows", "Hospital departments (OPD, Diagnostic Lab, Pharmacy) operate on disparate systems or manual paper logs, causing lost records and administrative delays."),
        ("Diagnostic Custody Gaps & Delay", "Diagnostic lab requisitions lack real-time chain-of-custody tracking. Critical/panic values often fail to alert prescribing physicians immediately, endangering patients."),
        ("Prescription Fulfillment Inefficiencies", "Paper and unverified prescriptions lead to medication transcription errors, lack of drug-drug interaction checks, and zero visibility into pharmacy dispensing status."),
        ("Absence of Patient-Centric Portals", "Patients struggle to access structured medical histories, verified diagnostic test results, active medication schedules, or accessible evidence-based guidance.")
    ]

    for i, (title, desc) in enumerate(problems):
        y = Inches(1.65 + i * 1.25)
        add_card(s3, Inches(0.8), y, Inches(11.733), Inches(1.1), WHITE, SLATE_300)
        
        b = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), y + Inches(0.18), Inches(0.4), Inches(0.74))
        b.fill.solid()
        b.fill.fore_color.rgb = WARNING_AMBER
        b.line.fill.background()

        tb = s3.shapes.add_textbox(Inches(1.55), y + Inches(0.12), Inches(10.8), Inches(0.85))
        tf_p = tb.text_frame
        tf_p.word_wrap = True
        p1 = tf_p.paragraphs[0]
        p1.text = f"Challenge {i+1}: {title}"
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_900

        p2 = tf_p.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = SLATE_600
        p2.space_before = Pt(3)

    # =========================================================================
    # SLIDE 4: Proposed Solution
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_header(s4, "Proposed Solution: The CareAI Platform", "Platform Concept", 4)

    flow_steps = [
        ("1. Patient Portal", "Registration, doctor search, appointment booking & report access.", PRIMARY_BLUE),
        ("2. Doctor Workspace", "Consultation queue, digital Rx drafting, AI drug safety & lab orders.", ACCENT_INDIGO),
        ("3. Lab Lifecycle", "Sample collection, analytical testing, panic alerts & verification.", PURPLE_600),
        ("4. Pharmacy Hub", "Prescription review, packaging status & secure dispensation.", SUCCESS_GREEN)
    ]
    for i, (title, sub, col) in enumerate(flow_steps):
        x = Inches(0.8 + i * 2.98)
        add_card(s4, x, Inches(1.65), Inches(2.8), Inches(1.8), WHITE, SLATE_300)
        
        strip = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.65), Inches(2.8), Inches(0.08))
        strip.fill.solid()
        strip.fill.fore_color.rgb = col
        strip.line.fill.background()

        tb = s4.shapes.add_textbox(x + Inches(0.15), Inches(1.85), Inches(2.5), Inches(1.5))
        tf_s = tb.text_frame
        tf_s.word_wrap = True
        p1 = tf_s.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = col
        
        p2 = tf_s.add_paragraph()
        p2.text = sub
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10)
        p2.font.color.rgb = SLATE_600
        p2.space_before = Pt(4)

    add_card(s4, Inches(0.8), Inches(3.7), Inches(11.733), Inches(3.0), SLATE_100, SLATE_300)
    sol_tb = s4.shapes.add_textbox(Inches(1.1), Inches(3.9), Inches(11.1), Inches(2.6))
    tf_sol = sol_tb.text_frame
    tf_sol.word_wrap = True
    
    p_hdr = tf_sol.paragraphs[0]
    p_hdr.text = "Key Architectural Value Propositions of CareAI:"
    p_hdr.font.name = FONT_HEADING
    p_hdr.font.size = Pt(14)
    p_hdr.font.bold = True
    p_hdr.font.color.rgb = NAVY_900

    points = [
        ("Single Source of Truth:", "All appointments, prescriptions, diagnostic results, and medical records stored in unified relational schema."),
        ("Strict Role-Based Access Control:", "5 distinct roles with hard isolation; Admins barred from viewing private patient medical documents (PHI barrier)."),
        ("Embedded Clinical AI Safety:", "Multi-drug interaction analysis, symptom triage with emergency warnings, and graceful offline fallback."),
        ("Production Cloud Readiness:", "Decoupled React Single Page Application on Vercel backed by FastAPI & PostgreSQL on Render.")
    ]
    for bold_prefix, text in points:
        p = tf_sol.add_paragraph()
        p.text = f"• {bold_prefix} {text}"
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = SLATE_600
        p.space_before = Pt(6)

    # =========================================================================
    # SLIDE 5: Key Objectives
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_header(s5, "Project Objectives & Technical Goals", "Target Deliverables", 5)

    objectives = [
        ("1. Multi-Role Healthcare Digitization", "Build a responsive, modern web application supporting 5 distinct user roles (Patient, Doctor, Admin, Lab Tech, Pharmacy Staff) with dedicated workflows."),
        ("2. End-to-End Clinical Lifecycle", "Implement seamless transitions between appointments, consultations, electronic prescription issuance, laboratory testing, and pharmacy dispensation."),
        ("3. Diagnostic Custody & Panic Alerting", "Design a 6-stage laboratory state machine tracking specimens from collection to release, with automated critical panic value detection and doctor alerts."),
        ("4. Pharmacy Dispensary Workflow", "Enable pharmacists to review prescriptions, track packaging progress, and confirm dispensation without modifying physician medical orders."),
        ("5. Guardrailed AI Clinical Assistance", "Integrate Google Gemini 1.5 Flash for drug safety and patient inquiries, fortified with emergency symptom triage and HTTP 503 fallback."),
        ("6. Enterprise Security & Automated Testing", "Implement Bcrypt password hashing, JWT stateless authentication, Alembic database migrations, and 165+ passing automated Pytest test cases.")
    ]

    for i, (title, desc) in enumerate(objectives):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.0)
        y = Inches(1.65 + row * 1.7)
        add_card(s5, x, y, Inches(5.733), Inches(1.5), WHITE, SLATE_300)

        tb = s5.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(5.3), Inches(1.2))
        tf_o = tb.text_frame
        tf_o.word_wrap = True
        p1 = tf_o.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(12.5)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_DARK

        p2 = tf_o.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = SLATE_600
        p2.space_before = Pt(3)

    # =========================================================================
    # SLIDE 6: System Users / 5-Role Architecture
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_slide_header(s6, "System Actors & 5-Role Access Architecture", "Role-Based Access Control", 6)

    roles = [
        ("PATIENT", "Primary Consumer", PRIMARY_BLUE, [
            "Search verified doctor directory",
            "Book & manage consultation slots",
            "Access active digital prescriptions",
            "View released diagnostic lab reports",
            "Interact with CareAI Health Assistant"
        ]),
        ("DOCTOR", "Clinical Provider", ACCENT_INDIGO, [
            "Configure weekly availability schedule",
            "Manage patient consultation queue",
            "Draft multi-drug digital prescriptions",
            "Execute AI drug-drug safety checks",
            "Order STAT & routine diagnostic panels"
        ]),
        ("ADMIN", "Governance & Vetting", NAVY_900, [
            "Review & approve doctor credentials",
            "Provision Lab Tech & Pharmacy staff",
            "Maintain standardized test catalog",
            "Monitor platform usage & stats",
            "Blocked from patient PHI docs (Privacy)"
        ]),
        ("LAB TECHNICIAN", "Diagnostic Specialist", PURPLE_600, [
            "Accession diagnostic requisition queue",
            "Collect & verify specimen integrity",
            "Enter quantitative analyte test values",
            "Flag critical/panic threshold values",
            "Verify & release official diagnostic reports"
        ]),
        ("PHARMACY STAFF", "Dispensary Pharmacist", SUCCESS_GREEN, [
            "Review incoming physician prescriptions",
            "Inspect dosage & drug safety summaries",
            "Advance order: Under Review → Ready",
            "Confirm dispensation & notify patient",
            "Protected: Cannot alter prescription orders"
        ])
    ]

    for i, (role_name, role_sub, col, duties) in enumerate(roles):
        x = Inches(0.8 + i * 2.38)
        add_card(s6, x, Inches(1.65), Inches(2.26), Inches(5.1), WHITE, SLATE_300)
        
        hdr = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.65), Inches(2.26), Inches(0.8))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = col
        hdr.line.fill.background()
        tf_h = hdr.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = role_name
        p_h.font.name = FONT_HEADING
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = WHITE
        p_h.alignment = PP_ALIGN.CENTER
        
        p_sub = tf_h.add_paragraph()
        p_sub.text = role_sub
        p_sub.font.name = FONT_BODY
        p_sub.font.size = Pt(8.5)
        p_sub.font.color.rgb = SLATE_300
        p_sub.alignment = PP_ALIGN.CENTER

        tb = s6.shapes.add_textbox(x + Inches(0.12), Inches(2.55), Inches(2.02), Inches(4.1))
        tf_r = tb.text_frame
        tf_r.word_wrap = True
        for j, duty in enumerate(duties):
            p = tf_r.paragraphs[0] if j == 0 else tf_r.add_paragraph()
            p.text = f"• {duty}"
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = SLATE_600
            if j > 0:
                p.space_before = Pt(6)

    # =========================================================================
    # SLIDE 7: High-Level System Architecture
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_slide_header(s7, "High-Level 3-Tier System Architecture", "System Design", 7)

    layers = [
        ("Tier 1: Presentation Layer (Client)", "React 18/19 SPA  •  Vite 5.4  •  React Router  •  Custom Glassmorphic CSS  •  Lucide Icons",
         ["Component-Driven UI with scoped sub-second HMR", "Role-Based Client Route Guards (<ProtectedRoute>)", "Global AuthContext with automatic JWT handling", "Responsive interfaces across Desktop and Mobile viewports"], PRIMARY_BLUE),
        ("Tier 2: Application & Gateway Layer (FastAPI)", "Python 3.12  •  FastAPI ASGI  •  Pydantic v2  •  OAuth2 / JWT  •  Passlib (Bcrypt)",
         ["Asynchronous high-throughput REST API controllers", "Strict Dependency Injection for DB sessions & RBAC roles", "Domain Service Layer encapsulating transaction boundaries", "Standardized OpenAPI 3.0 auto-generated documentation"], ACCENT_INDIGO),
        ("Tier 3: Persistence & External AI Layer", "PostgreSQL  •  SQLAlchemy 2.0 ORM  •  Alembic Migrations  •  Google Gemini 1.5 Flash",
         ["Relational persistence with ACID guarantees & foreign keys", "Alembic code-versioned database migrations to HEAD", "Abstracted AI Service Layer with symptom triage filtering", "HTTP 503 graceful degradation when AI service is offline"], SUCCESS_GREEN)
    ]

    for i, (title, stack, items, col) in enumerate(layers):
        y = Inches(1.65 + i * 1.7)
        add_card(s7, Inches(0.8), y, Inches(11.733), Inches(1.55), WHITE, SLATE_300)
        
        bar = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.2), Inches(1.55))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()

        tb = s7.shapes.add_textbox(Inches(1.15), y + Inches(0.1), Inches(11.2), Inches(1.35))
        tf_l = tb.text_frame
        tf_l.word_wrap = True
        
        p_t = tf_l.paragraphs[0]
        p_t.text = title
        p_t.font.name = FONT_HEADING
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_900

        p_s = tf_l.add_paragraph()
        p_s.text = f"Stack: {stack}"
        p_s.font.name = FONT_HEADING
        p_s.font.size = Pt(9.5)
        p_s.font.bold = True
        p_s.font.color.rgb = col
        p_s.space_before = Pt(2)

        bullets_text = "   |   ".join([f"• {it}" for it in items])
        p_b = tf_l.add_paragraph()
        p_b.text = bullets_text
        p_b.font.name = FONT_BODY
        p_b.font.size = Pt(9.5)
        p_b.font.color.rgb = SLATE_600
        p_b.space_before = Pt(4)

    # =========================================================================
    # SLIDE 8: Technology Stack Justification
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_slide_header(s8, "Technology Stack & Technical Justification", "Technology Selection", 8)

    tech_items = [
        ("Frontend: React + Vite", "Component reusability, declarative state, and instant development HMR with Rollup production tree-shaking."),
        ("Backend: FastAPI + Python", "High-performance asynchronous execution (ASGI), automatic OpenAPI schemas, and seamless Pydantic validation."),
        ("ORM: SQLAlchemy 2.0", "Enterprise query builder supporting explicit relationship modeling, eager/lazy loading, and database dialect portability."),
        ("Migrations: Alembic", "Deterministic, code-managed schema versioning allowing reproducible database upgrades across local, staging, and production."),
        ("Database: PostgreSQL", "ACID-compliant relational persistence, check constraints, foreign keys, and native JSON support for clinical metadata."),
        ("Security: JWT + Bcrypt", "Stateless authentication eliminating session lookup latency with slow 12-round salted password hashing."),
        ("AI: Google Gemini 1.5 Flash", "Large context window, multimodal clinical document understanding, fast inference, and medical terminology accuracy."),
        ("Deployment: Vercel + Render", "Edge CDN distribution for React assets and managed containerized cloud hosting for FastAPI and PostgreSQL.")
    ]

    for i, (tech, just) in enumerate(tech_items):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.0)
        y = Inches(1.65 + row * 1.28)
        add_card(s8, x, y, Inches(5.733), Inches(1.15), WHITE, SLATE_300)

        tb = s8.shapes.add_textbox(x + Inches(0.2), y + Inches(0.12), Inches(5.3), Inches(0.9))
        tf_tj = tb.text_frame
        tf_tj.word_wrap = True
        p1 = tf_tj.paragraphs[0]
        p1.text = f"✔  {tech}"
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_DARK

        p2 = tf_tj.add_paragraph()
        p2.text = just
        p2.font.name = FONT_BODY
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = SLATE_600
        p2.space_before = Pt(2)

    # =========================================================================
    # SLIDE 9: Patient & Doctor Workflow
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    add_slide_header(s9, "Patient & Doctor Clinical Consultation Journey", "Clinical Workflow", 9)

    add_card(s9, Inches(0.8), Inches(1.65), Inches(5.7), Inches(5.1), WHITE, SLATE_300)
    tb_p = s9.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.7))
    tf_pf = tb_p.text_frame
    tf_pf.word_wrap = True
    p_ph = tf_pf.paragraphs[0]
    p_ph.text = "Patient Experience Pipeline:"
    p_ph.font.name = FONT_HEADING
    p_ph.font.size = Pt(14)
    p_ph.font.bold = True
    p_ph.font.color.rgb = PRIMARY_BLUE

    patient_steps = [
        ("1. Discovery & Search", "Search approved doctor directory by specialty, fee, experience."),
        ("2. Slot Booking", "Select available doctor schedule slot and provide chief complaint."),
        ("3. Consultation Attended", "Doctor reviews medical profile and conducts consultation."),
        ("4. Digital Prescription", "View prescribed medications, dosage frequency & instructions."),
        ("5. Diagnostic Reports", "Access verified & released lab reports directly on portal."),
        ("6. CareAI Interaction", "Query AI Health Assistant for post-consultation health info.")
    ]
    for title, desc in patient_steps:
        p = tf_pf.add_paragraph()
        p.text = f"→ {title}: {desc}"
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.color.rgb = SLATE_600
        p.space_before = Pt(5)

    add_card(s9, Inches(6.833), Inches(1.65), Inches(5.7), Inches(5.1), WHITE, SLATE_300)
    tb_d = s9.shapes.add_textbox(Inches(7.033), Inches(1.8), Inches(5.3), Inches(4.7))
    tf_df = tb_d.text_frame
    tf_df.word_wrap = True
    p_dh = tf_df.paragraphs[0]
    p_dh.text = "Doctor Clinical Operations:"
    p_dh.font.name = FONT_HEADING
    p_dh.font.size = Pt(14)
    p_dh.font.bold = True
    p_dh.font.color.rgb = ACCENT_INDIGO

    doctor_steps = [
        ("1. Credential Vetting", "Register medical license & bio; vetted & approved by Admin."),
        ("2. Availability Scheduling", "Set recurring weekly consultation windows and slot duration."),
        ("3. Patient Queue", "Manage incoming appointments (Confirm / Complete / Cancel)."),
        ("4. AI Drug Safety Check", "Analyze drug-drug and drug-food interactions before issuing Rx."),
        ("5. Electronic Prescription", "Issue signed multi-medication digital prescription to patient & pharmacy."),
        ("6. Diagnostic Requisition", "Order STAT or routine laboratory test panels for technician execution.")
    ]
    for title, desc in doctor_steps:
        p = tf_df.add_paragraph()
        p.text = f"→ {title}: {desc}"
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.color.rgb = SLATE_600
        p.space_before = Pt(5)

    # =========================================================================
    # SLIDE 10: Lab Management Workflow
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    add_slide_header(s10, "Diagnostic Laboratory Lifecycle Management", "Diagnostic Workflow", 10)

    lab_stages = [
        ("1. SAMPLE_PENDING", "Doctor issues order requisition (STAT / Urgent / Routine). Test items accessioned in lab queue.", SLATE_500),
        ("2. SAMPLE_COLLECTED", "Phlebotomist collects specimen. Rejection check for compromised draws (Hemolyzed/Clotted).", PRIMARY_BLUE),
        ("3. IN_PROGRESS", "Specimen placed on analytical instruments. Technologist initiates biochemical testing.", ACCENT_INDIGO),
        ("4. RESULTS_ENTERED", "Quantitative analyte values recorded. Automated panic threshold alert checks triggered.", PURPLE_600),
        ("5. VERIFIED", "Supervising technician audits findings against QC reference intervals and signs off.", WARNING_AMBER),
        ("6. RELEASED", "Diagnostic report released to Patient Portal and ordering Doctor with immutable audit trail.", SUCCESS_GREEN)
    ]

    for i, (stage, desc, col) in enumerate(lab_stages):
        x = Inches(0.8 + (i % 3) * 3.98)
        y = Inches(1.65 + (i // 3) * 2.5)
        add_card(s10, x, y, Inches(3.75), Inches(2.3), WHITE, SLATE_300)

        tag = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), y + Inches(0.15), Inches(3.45), Inches(0.45))
        tag.fill.solid()
        tag.fill.fore_color.rgb = col
        tag.line.fill.background()
        tf_tg = tag.text_frame
        p_tg = tf_tg.paragraphs[0]
        p_tg.text = stage
        p_tg.font.name = FONT_HEADING
        p_tg.font.size = Pt(11)
        p_tg.font.bold = True
        p_tg.font.color.rgb = WHITE
        p_tg.alignment = PP_ALIGN.CENTER

        tb = s10.shapes.add_textbox(x + Inches(0.15), y + Inches(0.7), Inches(3.45), Inches(1.4))
        tf_d = tb.text_frame
        tf_d.word_wrap = True
        p = tf_d.paragraphs[0]
        p.text = desc
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.color.rgb = SLATE_600

    # =========================================================================
    # SLIDE 11: Pharmacy Workflow
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    add_slide_header(s11, "Pharmacy Dispensary & Fulfillment Workflow", "Pharmacy Workflow", 11)

    rx_stages = [
        ("PRESCRIBED", "Doctor issues digital prescription. Order enters dispensary queue automatically.", PRIMARY_BLUE),
        ("UNDER_REVIEW", "Pharmacist audits prescription items, dosage guidelines, and AI safety warnings.", ACCENT_INDIGO),
        ("READY FOR PICKUP", "Medications packaged; automated in-app pickup notification dispatched to Patient.", PURPLE_600),
        ("DISPENSED", "Patient receives medication. Pharmacist signs off; dispensation timestamp permanently recorded.", SUCCESS_GREEN)
    ]

    for i, (stg_name, stg_desc, col) in enumerate(rx_stages):
        x = Inches(0.8 + i * 2.98)
        add_card(s11, x, Inches(1.65), Inches(2.2), Inches(2.2), WHITE, SLATE_300)
        
        hdr = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), Inches(1.8), Inches(2.5), Inches(0.5))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = col
        hdr.line.fill.background()
        tf_h = hdr.text_frame
        p = tf_h.paragraphs[0]
        p.text = f"Stage {i+1}: {stg_name}"
        p.font.name = FONT_HEADING
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        tb = s11.shapes.add_textbox(x + Inches(0.15), Inches(2.4), Inches(2.5), Inches(1.3))
        tf_s = tb.text_frame
        tf_s.word_wrap = True
        p_s = tf_s.paragraphs[0]
        p_s.text = stg_desc
        p_s.font.name = FONT_BODY
        p_s.font.size = Pt(10)
        p_s.font.color.rgb = SLATE_600

    add_card(s11, Inches(0.8), Inches(4.1), Inches(11.733), Inches(2.6), SLATE_100, SLATE_300)
    tb_call = s11.shapes.add_textbox(Inches(1.1), Inches(4.3), Inches(11.1), Inches(2.2))
    tf_c = tb_call.text_frame
    tf_c.word_wrap = True
    p_ch = tf_c.paragraphs[0]
    p_ch.text = "Key Clinical Safety & Security Guardrails in Pharmacy Module:"
    p_ch.font.name = FONT_HEADING
    p_ch.font.size = Pt(13)
    p_ch.font.bold = True
    p_ch.font.color.rgb = NAVY_900

    pharm_points = [
        ("Write-Isolation Barrier:", "Pharmacists CANNOT modify prescribed drug names, dosages, or physician diagnoses. They can only record pharmacist dispensing notes and advance fulfillment status."),
        ("Clinical Safety Visibility:", "Pharmacists inspect AI drug-drug and drug-food interaction reports generated during physician prescribing."),
        ("Multi-Role Notification Triggers:", "Advancing to READY alerts the Patient; confirming DISPENSED notifies both Patient and Prescribing Doctor.")
    ]
    for b_title, b_text in pharm_points:
        p = tf_c.add_paragraph()
        p.text = f"• {b_title} {b_text}"
        p.font.name = FONT_BODY
        p.font.size = Pt(10.5)
        p.font.color.rgb = SLATE_600
        p.space_before = Pt(4)

    # =========================================================================
    # SLIDE 12: AI Integration Architecture
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    add_slide_header(s12, "Responsible AI Architecture & Clinical Guardrails", "Artificial Intelligence", 12)

    ai_steps = [
        ("Layer 1: Deterministic Emergency Triage", "Incoming patient prompts scanned with regex/keyword rules for acute emergencies (chest pain, anaphylaxis, severe dyspnea). Instantly returns emergency helpline alerts without model latency.", PRIMARY_BLUE),
        ("Layer 2: System Prompt & PHI Guardrails", "Role-aware prompt engineering constrains LLM output to evidence-based medical definitions, attaches legal medical disclaimers, and prevents PHI data leakage.", ACCENT_INDIGO),
        ("Layer 3: Google Gemini 1.5 Flash API", "Processes multi-turn health inquiries, multi-medication drug-drug interaction matrices, and patient medical document OCR analysis with high clinical terminology accuracy.", PURPLE_600),
        ("Layer 4: Graceful Fault Tolerance (HTTP 503)", "If Gemini API is unreachable or rate-limited, system catches AIProviderUnavailableError and returns HTTP 503 with structured retry advice. Core clinical EHR remains 100% operational.", SUCCESS_GREEN)
    ]

    for i, (title, desc, col) in enumerate(ai_steps):
        y = Inches(1.65 + i * 1.28)
        add_card(s12, Inches(0.8), y, Inches(11.733), Inches(1.15), WHITE, SLATE_300)

        bar = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.18), Inches(1.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()

        tb = s12.shapes.add_textbox(Inches(1.15), y + Inches(0.1), Inches(11.2), Inches(0.95))
        tf_ai = tb.text_frame
        tf_ai.word_wrap = True
        p1 = tf_ai.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = col

        p2 = tf_ai.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10)
        p2.font.color.rgb = SLATE_600
        p2.space_before = Pt(2)

    # =========================================================================
    # SLIDE 13: Security & RBAC
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    add_slide_header(s13, "Security, Authentication & Data Privacy", "Security & RBAC", 13)

    sec_cards = [
        ("Cryptographic Password Hashing", "Passlib with Bcrypt algorithm (12 salt rounds) prevents rainbow table and GPU brute-force attacks. Passwords never stored in plaintext.", PRIMARY_BLUE),
        ("Stateless JWT Authentication", "Tokens signed with HMAC-SHA256 (HS256) containing sub (email), user ID, and role claims. Eliminates database session lookup latency.", ACCENT_INDIGO),
        ("Role-Based Access Control (RBAC)", "FastAPI dependency guards (require_role, require_roles) enforce strict access boundaries at the API gateway layer.", PURPLE_600),
        ("Protected Health Info (PHI) Barrier", "Admins manage platform users but are strictly forbidden (HTTP 403) from accessing or downloading patient medical documents.", SUCCESS_GREEN)
    ]

    for i, (title, desc, col) in enumerate(sec_cards):
        col_idx = i % 2
        row_idx = i // 2
        x = Inches(0.8 + col_idx * 6.0)
        y = Inches(1.65 + row_idx * 2.5)
        add_card(s13, x, y, Inches(5.733), Inches(2.3), WHITE, SLATE_300)

        strip = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.2), Inches(5.333), Inches(0.45))
        strip.fill.solid()
        strip.fill.fore_color.rgb = SLATE_100
        strip.line.fill.background()
        tf_s = strip.text_frame
        p_st = tf_s.paragraphs[0]
        p_st.text = f"🔒  {title}"
        p_st.font.name = FONT_HEADING
        p_st.font.size = Pt(12)
        p_st.font.bold = True
        p_st.font.color.rgb = col

        tb = s13.shapes.add_textbox(x + Inches(0.2), y + Inches(0.75), Inches(5.333), Inches(1.4))
        tf_d = tb.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = SLATE_600

    # =========================================================================
    # SLIDE 14: Database Design & Schema
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    add_slide_header(s14, "Relational Database Design & Entity Architecture", "Database Schema", 14)

    db_groups = [
        ("Identity & Roles", ["users (Base credentials & roles)", "patient_profiles (Clinical demographics)", "doctor_profiles (Specialty & approvals)"], PRIMARY_BLUE),
        ("Consultations & Scheduling", ["doctor_availabilities (Weekly slots)", "doctor_unavailable_dates (Vacations)", "appointments (Booking lifecycle)"], ACCENT_INDIGO),
        ("Prescriptions & Safety", ["prescriptions (Status & diagnosis)", "prescription_items (Drug dosage/route)", "ai_reports (Drug interaction findings)"], PURPLE_600),
        ("Laboratory Management", ["lab_tests (Standard catalog)", "lab_orders & lab_order_items", "lab_samples, lab_results, lab_audit_events"], SUCCESS_GREEN),
        ("Documents & Intelligence", ["medical_documents (Patient uploads)", "medical_document_analyses (OCR summaries)", "ai_conversations & ai_messages"], WARNING_AMBER),
        ("Platform Communications", ["notifications (In-app multi-role alerts)", "Integrity: Foreign Keys & Cascades", "Alembic Version-Controlled Migrations"], NAVY_900)
    ]

    for i, (grp_name, tbls, col) in enumerate(db_groups):
        col_idx = i % 3
        row_idx = i // 3
        x = Inches(0.8 + col_idx * 3.98)
        y = Inches(1.65 + row_idx * 2.55)
        add_card(s14, x, y, Inches(3.75), Inches(2.35), WHITE, SLATE_300)

        hdr = s14.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.75), Inches(0.45))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = col
        hdr.line.fill.background()
        tf_h = hdr.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = grp_name
        p_h.font.name = FONT_HEADING
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = WHITE
        p_h.alignment = PP_ALIGN.CENTER

        tb = s14.shapes.add_textbox(x + Inches(0.15), y + Inches(0.55), Inches(3.45), Inches(1.7))
        tf_t = tb.text_frame
        tf_t.word_wrap = True
        for j, tbl in enumerate(tbls):
            p = tf_t.paragraphs[0] if j == 0 else tf_t.add_paragraph()
            p.text = f"• {tbl}"
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = SLATE_600
            if j > 0:
                p.space_before = Pt(4)

    # =========================================================================
    # SLIDE 15: Key Features Grid
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    add_slide_header(s15, "CareAI Core Platform Feature Matrix", "Feature Summary", 15)

    features = [
        ("Multi-Role Authentication", "Seamless login, JWT issuance, and auto-redirection to role-tailored dashboards."),
        ("Doctor Vetting & Approval", "Admin approval gate before doctor profiles become bookable in public directory."),
        ("Dynamic Slot Scheduling", "Recurring weekly availability and ad-hoc unavailable date blocking."),
        ("Digital Prescription Issuance", "Structured multi-item prescriptions with dosage, frequency, and instructions."),
        ("AI Drug Interaction Screening", "Multi-drug and drug-food safety analysis with risk severity classifications."),
        ("Diagnostic Lab Custody Chain", "6-stage lab order lifecycle with specimen rejection and critical panic flags."),
        ("Pharmacy Fulfillment Tracking", "Prescription status pipeline from review to ready-for-pickup and dispensation."),
        ("Medical Document Vault", "Secure upload of patient documents with automated OCR and clinical summaries."),
        ("AI Health Assistant", "Context-aware conversational assistant with emergency symptom triage filters.")
    ]

    for i, (title, desc) in enumerate(features):
        col_idx = i % 3
        row_idx = i // 3
        x = Inches(0.8 + col_idx * 3.98)
        y = Inches(1.65 + row_idx * 1.7)
        add_card(s15, x, y, Inches(3.75), Inches(1.55), WHITE, SLATE_300)

        tb = s15.shapes.add_textbox(x + Inches(0.15), y + Inches(0.12), Inches(3.45), Inches(1.3))
        tf_f = tb.text_frame
        tf_f.word_wrap = True
        p1 = tf_f.paragraphs[0]
        p1.text = f"✔ {title}"
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_DARK

        p2 = tf_f.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(9)
        p2.font.color.rgb = SLATE_600
        p2.space_before = Pt(2)

    # =========================================================================
    # SLIDE 16: Testing & Quality Assurance
    # =========================================================================
    s16 = prs.slides.add_slide(blank_layout)
    add_slide_header(s16, "Automated Testing & Quality Assurance", "Quality Engineering", 16)

    add_card(s16, Inches(0.8), Inches(1.65), Inches(4.5), Inches(5.1), SUCCESS_BG, SUCCESS_GREEN)
    tb_m = s16.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(4.1), Inches(4.7))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    
    p_mh = tf_m.paragraphs[0]
    p_mh.text = "165 / 165"
    p_mh.font.name = FONT_HEADING
    p_mh.font.size = Pt(44)
    p_mh.font.bold = True
    p_mh.font.color.rgb = SUCCESS_GREEN
    p_mh.alignment = PP_ALIGN.CENTER

    p_ms = tf_m.add_paragraph()
    p_ms.text = "AUTOMATED PYTEST TESTS PASSED"
    p_ms.font.name = FONT_HEADING
    p_ms.font.size = Pt(12)
    p_ms.font.bold = True
    p_ms.font.color.rgb = NAVY_900
    p_ms.alignment = PP_ALIGN.CENTER
    p_ms.space_before = Pt(4)

    p_mb = tf_m.add_paragraph()
    p_mb.text = "100% Pass Rate across all unit, integration, RBAC security, and cross-role workflow test suites in ~6.8 seconds."
    p_mb.font.name = FONT_BODY
    p_mb.font.size = Pt(10.5)
    p_mb.font.color.rgb = SLATE_600
    p_mb.alignment = PP_ALIGN.CENTER
    p_mb.space_before = Pt(10)

    test_suites = [
        ("Auth & RBAC Isolation Tests", "Verifies JWT issuance, password hashing, and role access barriers (e.g. Admin PHI blocking)."),
        ("Clinical Appointment & Availability", "Tests doctor weekly availability slot generation, conflict detection, and booking state transitions."),
        ("Prescription & Pharmacy Lifecycle", "Verifies multi-item Rx drafting, AI report attachment, and dispensary progression."),
        ("Lab Diagnostic Chain-of-Custody", "Validates sample collection, rejection criteria, panic threshold flagging, and verification signoff."),
        ("AI Safety Triage & Fault Handling", "Tests deterministic emergency triage triggers and graceful HTTP 503 fallback handling.")
    ]

    for i, (suite_title, suite_desc) in enumerate(test_suites):
        y = Inches(1.65 + i * 1.02)
        add_card(s16, Inches(5.6), y, Inches(6.933), Inches(0.92), WHITE, SLATE_300)

        tb_s = s16.shapes.add_textbox(Inches(5.8), y + Inches(0.08), Inches(6.5), Inches(0.75))
        tf_st = tb_s.text_frame
        tf_st.word_wrap = True
        p1 = tf_st.paragraphs[0]
        p1.text = f"✔ {suite_title}"
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_DARK

        p2 = tf_st.add_paragraph()
        p2.text = suite_desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(9)
        p2.font.color.rgb = SLATE_600
        p2.space_before = Pt(1)

    # =========================================================================
    # SLIDE 17: Deployment Architecture
    # =========================================================================
    s17 = prs.slides.add_slide(blank_layout)
    add_slide_header(s17, "Production Cloud Deployment Architecture", "DevOps & Deployment", 17)

    dep_cards = [
        ("Frontend Deployment: Vercel", "React 18 / Vite SPA  •  Global Edge CDN", [
            "Continuous deployment from GitHub main branch",
            "Automatic SSL/TLS termination and HTTP/2 acceleration",
            "Sub-second asset distribution across edge nodes",
            "SPA rewrites for react-router-dom deep linking"
        ], PRIMARY_BLUE),
        ("Backend Deployment: Render", "FastAPI ASGI  •  Python 3.12 Uvicorn", [
            "Containerized Python runtime with gunicorn/uvicorn workers",
            "Automated Alembic database migration execution on build",
            "Environment variable injection for secret key & Gemini API",
            "Automatic health check endpoint monitoring (/api/v1/health)"
        ], ACCENT_INDIGO),
        ("Database: Managed PostgreSQL", "Render Managed PostgreSQL  •  ACID Persistence", [
            "Production relational database with SSL connection strings",
            "Automated nightly backups and point-in-time recovery",
            "Full support for native JSON / JSONB data types",
            "Enforced Foreign Key constraints and indexing"
        ], SUCCESS_GREEN)
    ]

    for i, (title, sub, bullets, col) in enumerate(dep_cards):
        x = Inches(0.8 + i * 3.98)
        add_card(s17, x, Inches(1.65), Inches(3.75), Inches(5.1), WHITE, SLATE_300)

        hdr = s17.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.65), Inches(3.75), Inches(0.85))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = col
        hdr.line.fill.background()
        tf_h = hdr.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = title
        p_h.font.name = FONT_HEADING
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = WHITE
        p_h.alignment = PP_ALIGN.CENTER
        
        p_s = tf_h.add_paragraph()
        p_s.text = sub
        p_s.font.name = FONT_BODY
        p_s.font.size = Pt(8.5)
        p_s.font.color.rgb = SLATE_300
        p_s.alignment = PP_ALIGN.CENTER

        tb = s17.shapes.add_textbox(x + Inches(0.15), Inches(2.6), Inches(3.45), Inches(4.0))
        tf_b = tb.text_frame
        tf_b.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf_b.paragraphs[0] if j == 0 else tf_b.add_paragraph()
            p.text = f"• {b}"
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = SLATE_600
            if j > 0:
                p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 18: Project Results & Deliverables
    # =========================================================================
    s18 = prs.slides.add_slide(blank_layout)
    add_slide_header(s18, "Summary of Project Results & Deliverables", "Outcomes", 18)

    results = [
        ("Full-Stack 5-Role Platform", "Delivered complete healthcare SaaS with 5 dedicated user portals, responsive glassmorphic UI, and full REST API."),
        ("165/165 Automated Tests", "Achieved 100% test pass rate across auth, appointments, prescriptions, diagnostic workflows, and AI safety."),
        ("Automated Alembic Migrations", "Established clean, repeatable schema versioning across 21 relational database tables."),
        ("Responsible AI Implementation", "Implemented drug-drug interaction checks, symptom triage emergency filters, and graceful 503 offline fallback."),
        ("Production Cloud Deployment", "Deployed on Vercel (Frontend) and Render (Backend + PostgreSQL) with live HTTPS endpoints."),
        ("Verified Seed Demo Dataset", "Configured idempotent seed script creating realistic clinical scenarios across all 5 roles.")
    ]

    for i, (title, desc) in enumerate(results):
        col_idx = i % 2
        row_idx = i // 2
        x = Inches(0.8 + col_idx * 6.0)
        y = Inches(1.65 + row_idx * 1.7)
        add_card(s18, x, y, Inches(5.733), Inches(1.5), WHITE, SLATE_300)

        tb = s18.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(5.3), Inches(1.2))
        tf_r = tb.text_frame
        tf_r.word_wrap = True
        p1 = tf_r.paragraphs[0]
        p1.text = f"✔ {title}"
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(12.5)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_DARK

        p2 = tf_r.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10)
        p2.font.color.rgb = SLATE_600
        p2.space_before = Pt(3)

    # =========================================================================
    # SLIDE 19: Current System Limitations
    # =========================================================================
    s19 = prs.slides.add_slide(blank_layout)
    add_slide_header(s19, "Honest Analysis of Current System Limitations", "System Boundaries", 19)

    limits = [
        ("Clinical Decision Boundary", "CareAI is an assistive management and clinical support prototype. It is NOT a certified Class II Software as a Medical Device (SaMD) and does not replace certified physician diagnosis."),
        ("Notification Transport Mechanism", "Notifications are currently stored and queried via relational database polling. Real-time WebSockets, push notifications (FCM/APNs), and SMS/Email gateways (Twilio/SendGrid) are not yet integrated."),
        ("File Storage Architecture", "Document uploads are handled via local filesystem storage service. Enterprise production requires cloud object storage (AWS S3 / GCP GCS) with presigned URLs and malware scanning."),
        ("Pharmacy Live Stock Decrementing", "The pharmacy module tracks fulfillment lifecycle status progression, but does not maintain a live inventory count table that automatically decrements physical pill counts upon dispensing."),
        ("Database Tenancy Model", "Operates on a unified database with logical row-level tenant separation. Does not implement physical schema-per-hospital isolation for enterprise hospital networks.")
    ]

    for i, (title, desc) in enumerate(limits):
        y = Inches(1.65 + i * 1.02)
        add_card(s19, Inches(0.8), y, Inches(11.733), Inches(0.92), WHITE, SLATE_300)

        b = s19.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.18), Inches(0.92))
        b.fill.solid()
        b.fill.fore_color.rgb = WARNING_AMBER
        b.line.fill.background()

        tb = s19.shapes.add_textbox(Inches(1.15), y + Inches(0.08), Inches(11.2), Inches(0.75))
        tf_l = tb.text_frame
        tf_l.word_wrap = True
        p1 = tf_l.paragraphs[0]
        p1.text = f"• {title}"
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_900

        p2 = tf_l.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = SLATE_600
        p2.space_before = Pt(1)

    # =========================================================================
    # SLIDE 20: Future Scope & Roadmap
    # =========================================================================
    s20 = prs.slides.add_slide(blank_layout)
    add_slide_header(s20, "Future Scope & Production Roadmap", "Roadmap", 20)

    phases = [
        ("Phase 1: Real-Time Tele-Health", "WebSockets & WebRTC", PRIMARY_BLUE, [
            "Bidirectional WebSockets for instantaneous real-time notifications",
            "WebRTC peer-to-peer encrypted video rooms for virtual tele-consultations",
            "In-app real-time doctor-patient messaging"
        ]),
        ("Phase 2: Interoperability & Standards", "FHIR / HL7 & DICOM", ACCENT_INDIGO, [
            "Fast Healthcare Interoperability Resources (FHIR) standard JSON APIs",
            "HL7 interface for legacy Hospital Information System (HIS) integration",
            "Web-based DICOM medical radiographic image viewer (X-Ray/CT/MRI)"
        ]),
        ("Phase 3: Mobile & Inventory", "React Native & Barcodes", PURPLE_600, [
            "Cross-platform React Native iOS and Android mobile applications",
            "Live pharmacy stock inventory tracking with automatic decrementing",
            "Barcode / QR code medication scanning during dispensation"
        ]),
        ("Phase 4: Advanced Edge AI", "On-Premise LLM & Compliance", SUCCESS_GREEN, [
            "Embedded on-premise local clinical LLM (e.g. BioMistral via ONNX)",
            "Third-party HIPAA & GDPR security compliance certification audit",
            "Automated clinical coding (ICD-10 / CPT billing code recommendations)"
        ])
    ]

    for i, (title, sub, col, bullets) in enumerate(phases):
        x = Inches(0.8 + i * 2.98)
        add_card(s20, x, Inches(1.65), Inches(2.8), Inches(5.1), WHITE, SLATE_300)

        hdr = s20.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.65), Inches(2.8), Inches(0.8))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = col
        hdr.line.fill.background()
        tf_h = hdr.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = title
        p_h.font.name = FONT_HEADING
        p_h.font.size = Pt(10)
        p_h.font.bold = True
        p_h.font.color.rgb = WHITE
        p_h.alignment = PP_ALIGN.CENTER
        
        p_s = tf_h.add_paragraph()
        p_s.text = sub
        p_s.font.name = FONT_BODY
        p_s.font.size = Pt(8.5)
        p_s.font.color.rgb = SLATE_300
        p_s.alignment = PP_ALIGN.CENTER

        tb = s20.shapes.add_textbox(x + Inches(0.12), Inches(2.55), Inches(2.55), Inches(4.1))
        tf_b = tb.text_frame
        tf_b.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf_b.paragraphs[0] if j == 0 else tf_b.add_paragraph()
            p.text = f"• {b}"
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = SLATE_600
            if j > 0:
                p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 21: Conclusion
    # =========================================================================
    s21 = prs.slides.add_slide(blank_layout)
    add_slide_header(s21, "Project Conclusion & Technical Impact", "Conclusion", 21)

    add_card(s21, Inches(0.8), Inches(1.65), Inches(11.733), Inches(5.1), WHITE, SLATE_300)
    tb_c = s21.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(11.1), Inches(4.6))
    tf_con = tb_c.text_frame
    tf_con.word_wrap = True

    p_ch = tf_con.paragraphs[0]
    p_ch.text = "Key Takeaways from CareAI Development:"
    p_ch.font.name = FONT_HEADING
    p_ch.font.size = Pt(16)
    p_ch.font.bold = True
    p_ch.font.color.rgb = NAVY_900

    concl_points = [
        ("Unified Healthcare Experience:", "Successfully eliminated fragmented paper processes by combining appointments, digital prescriptions, laboratory chain-of-custody, and pharmacy fulfillment into a single reactive SaaS platform."),
        ("Strict Role-Based Security Architecture:", "Implemented enterprise RBAC across 5 roles with protected health information (PHI) barriers, ensuring data confidentiality and clinical custody compliance."),
        ("Responsible & Guardrailed AI:", "Demonstrated that Generative AI (Google Gemini) can be integrated responsibly into healthcare systems with deterministic emergency triage, prompt guardrails, and graceful offline fallback."),
        ("Rigorous Software Engineering:", "Achieved 100% automated test pass rate (165/165 tests) with full cloud deployment on Vercel and Render backed by PostgreSQL."),
        ("Solid Foundation for Future Expansion:", "Provides a modular, extensible architecture ready for real-time WebSockets, WebRTC video consultations, and FHIR/HL7 hospital interoperability.")
    ]
    for b_head, b_body in concl_points:
        p = tf_con.add_paragraph()
        p.text = f"✔ {b_head} {b_body}"
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = SLATE_600
        p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 22: Thank You / Q&A
    # =========================================================================
    s22 = prs.slides.add_slide(blank_layout)
    
    bg_end = s22.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg_end.fill.solid()
    bg_end.fill.fore_color.rgb = NAVY_900
    bg_end.line.fill.background()

    tb_end = s22.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(4.0))
    tf_e = tb_end.text_frame
    tf_e.word_wrap = True
    
    p_ty = tf_e.paragraphs[0]
    p_ty.text = "Thank You!"
    p_ty.font.name = FONT_HEADING
    p_ty.font.size = Pt(46)
    p_ty.font.bold = True
    p_ty.font.color.rgb = WHITE
    p_ty.alignment = PP_ALIGN.CENTER

    p_qa = tf_e.add_paragraph()
    p_qa.text = "Questions & Technical Discussion"
    p_qa.font.name = FONT_HEADING
    p_qa.font.size = Pt(22)
    p_qa.font.color.rgb = PRIMARY_BLUE
    p_qa.alignment = PP_ALIGN.CENTER
    p_qa.space_before = Pt(10)

    p_proj = tf_e.add_paragraph()
    p_proj.text = "CareAI: AI-Powered Healthcare Management SaaS Platform"
    p_proj.font.name = FONT_BODY
    p_proj.font.size = Pt(14)
    p_proj.font.color.rgb = SLATE_300
    p_proj.alignment = PP_ALIGN.CENTER
    p_proj.space_before = Pt(16)

    card_links = s22.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.6), Inches(4.6), Inches(8.1), Inches(1.5))
    card_links.fill.solid()
    card_links.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card_links.line.color.rgb = RGBColor(51, 65, 85)
    tf_cl = card_links.text_frame
    tf_cl.word_wrap = True
    
    p_gh = tf_cl.paragraphs[0]
    p_gh.text = "GitHub Repository: https://github.com/Niharika212006/CareAI-SaaS"
    p_gh.font.name = FONT_BODY
    p_gh.font.size = Pt(11)
    p_gh.font.color.rgb = WHITE
    p_gh.alignment = PP_ALIGN.CENTER

    p_dep = tf_cl.add_paragraph()
    p_dep.text = "Frontend: Vercel Deployment  •  Backend & DB: Render Cloud (FastAPI + PostgreSQL)"
    p_dep.font.name = FONT_BODY
    p_dep.font.size = Pt(10.5)
    p_dep.font.color.rgb = PRIMARY_BLUE
    p_dep.alignment = PP_ALIGN.CENTER
    p_dep.space_before = Pt(6)

    p_stat = tf_cl.add_paragraph()
    p_stat.text = "Automated Test Suite: 165/165 Passing  •  Role Fidelity: 5 User Roles Verified"
    p_stat.font.name = FONT_BODY
    p_stat.font.size = Pt(10)
    p_stat.font.color.rgb = SUCCESS_GREEN
    p_stat.alignment = PP_ALIGN.CENTER
    p_stat.space_before = Pt(4)

    # Save presentation
    prs.save(output_filepath)
    print(f"[SUCCESS] PowerPoint presentation saved successfully to: {output_filepath}")
    print(f"Total Slides Generated: {len(prs.slides)}")

if __name__ == "__main__":
    out_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "CareAI_Final_Year_Project_Presentation.pptx"))
    if len(sys.argv) > 1:
        out_path = sys.argv[1]
    create_presentation(out_path)
